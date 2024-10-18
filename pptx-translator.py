import argparse
import boto3
from botocore.exceptions import ClientError
from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID

LANGUAGE_CODE_TO_LANGUAGE_ID = {
"""
Dict that maps Amazon Translate language code to MSO_LANGUAGE_ID enum value.

- Amazon Translate language codes: https://docs.aws.amazon.com/translate/latest/dg/what-is.html#what-is-languages
- python-pptx MSO_LANGUAGE_ID enum: https://python-pptx.readthedocs.io/en/latest/api/enum/MsoLanguageId.html

python-pptx doesn't support:
    - Azerbaijani (az)
    - Persian (fa)
    - Dari (fa-AF)
    - Tagalog (tl)
"""
    'af': MSO_LANGUAGE_ID.AFRIKAANS,
    'am': MSO_LANGUAGE_ID.AMHARIC,
    'ar': MSO_LANGUAGE_ID.ARABIC,
    'bg': MSO_LANGUAGE_ID.BULGARIAN,
    'bn': MSO_LANGUAGE_ID.BENGALI,
    'bs': MSO_LANGUAGE_ID.BOSNIAN,
    'cs': MSO_LANGUAGE_ID.CZECH,
    'da': MSO_LANGUAGE_ID.DANISH,
    'de': MSO_LANGUAGE_ID.GERMAN,
    'el': MSO_LANGUAGE_ID.GREEK,
    'en': MSO_LANGUAGE_ID.ENGLISH_US,
    'es': MSO_LANGUAGE_ID.SPANISH,
    'et': MSO_LANGUAGE_ID.ESTONIAN,
    'fi': MSO_LANGUAGE_ID.FINNISH,
    'fr': MSO_LANGUAGE_ID.FRENCH,
    'fr-CA': MSO_LANGUAGE_ID.FRENCH_CANADIAN,
    'ha': MSO_LANGUAGE_ID.HAUSA,
    'he': MSO_LANGUAGE_ID.HEBREW,
    'hi': MSO_LANGUAGE_ID.HINDI,
    'hr': MSO_LANGUAGE_ID.CROATIAN,
    'hu': MSO_LANGUAGE_ID.HUNGARIAN,
    'id': MSO_LANGUAGE_ID.INDONESIAN,
    'it': MSO_LANGUAGE_ID.ITALIAN,
    'ja': MSO_LANGUAGE_ID.JAPANESE,
    'ka': MSO_LANGUAGE_ID.GEORGIAN,
    'ko': MSO_LANGUAGE_ID.KOREAN,
    'lv': MSO_LANGUAGE_ID.LATVIAN,
    'ms': MSO_LANGUAGE_ID.MALAYSIAN,
    'nl': MSO_LANGUAGE_ID.DUTCH,
    'no': MSO_LANGUAGE_ID.NORWEGIAN_BOKMOL,
    'pl': MSO_LANGUAGE_ID.POLISH,
    'ps': MSO_LANGUAGE_ID.PASHTO,
    'pt': MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE,
    'ro': MSO_LANGUAGE_ID.ROMANIAN,
    'ru': MSO_LANGUAGE_ID.RUSSIAN,
    'sk': MSO_LANGUAGE_ID.SLOVAK,
    'sl': MSO_LANGUAGE_ID.SLOVENIAN,
    'so': MSO_LANGUAGE_ID.SOMALI,
    'sq': MSO_LANGUAGE_ID.ALBANIAN,
    'sr': MSO_LANGUAGE_ID.SERBIAN_LATIN,
    'sv': MSO_LANGUAGE_ID.SWEDISH,
    'sw': MSO_LANGUAGE_ID.SWAHILI,
    'ta': MSO_LANGUAGE_ID.TAMIL,
    'th': MSO_LANGUAGE_ID.THAI,
    'tr': MSO_LANGUAGE_ID.TURKISH,
    'uk': MSO_LANGUAGE_ID.UKRAINIAN,
    'ur': MSO_LANGUAGE_ID.URDU,
    'vi': MSO_LANGUAGE_ID.VIETNAMESE,
    'zh': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    'zh-TW': MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
}

TERMINOLOGY_NAME = 'pptx-translator-terminology'

translate = boto3.client(service_name='translate')

def translate_presentation(presentation, source_language_code, target_language_code, terminology_names):
    for slide_index, slide in enumerate(presentation.slides, start=1):
        print(f'Slide {slide_index} of {len(presentation.slides)}')
        
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        translate_text_frame(cell.text_frame, source_language_code, target_language_code, terminology_names)
            elif shape.has_text_frame:
                translate_text_frame(shape.text_frame, source_language_code, target_language_code, terminology_names)
        
        if slide.has_notes_slide:
            translate_text_frame(slide.notes_slide.notes_text_frame, source_language_code, target_language_code, terminology_names)

def translate_text_frame(text_frame, source_language_code, target_language_code, terminology_names):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                try:
                    response = translate.translate_text(
                        Text=run.text,
                        SourceLanguageCode=source_language_code,
                        TargetLanguageCode=target_language_code,
                        TerminologyNames=terminology_names
                    )
                    # original text if translation fails
                    run.text = response.get('TranslatedText', run.text)  
                except ClientError as client_error:
                    if client_error.response['Error']['Code'] == 'ValidationException':
                        print('Invalid text. Ignoring...')

def import_terminology(terminology_file_path):
    print(f'Importing terminology data from {terminology_file_path}...')
    with open(terminology_file_path, 'rb') as f:
        translate.import_terminology(Name=TERMINOLOGY_NAME,
                                     MergeStrategy='OVERWRITE',
                                     TerminologyData={'File': bytearray(f.read()), 'Format': 'CSV'})

def main():
    argument_parser = argparse.ArgumentParser(
            'Translates pptx files from source language to target language using Amazon Translate service')
    argument_parser.add_argument(
            'source_language_code', type=str,
            help='The language code for the language of the source text. Example: en')
    argument_parser.add_argument(
            'target_language_code', type=str,
            help='The language code requested for the language of the target text. Example: pt')
    argument_parser.add_argument(
            'input_file_path', type=str,
            help='The path of the pptx file that should be translated')
    argument_parser.add_argument(
            '--terminology', type=str,
            help='The path of the terminology CSV file')
    args = argument_parser.parse_args()

    terminology_names = []
    if args.terminology:
        import_terminology(args.terminology)
        terminology_names = [TERMINOLOGY_NAME]

    print(f'Translating {args.input_file_path} from {args.source_language_code} to {args.target_language_code}...')
    presentation = Presentation(args.input_file_path)
    translate_presentation(presentation,
                           args.source_language_code,
                           args.target_language_code,
                           terminology_names)

    output_file_path = args.input_file_path.replace(
            '.pptx', f'-{args.target_language_code}.pptx')
    print(f'Saving {output_file_path}...')
    presentation.save(output_file_path)

if __name__ == '__main__':
    main()